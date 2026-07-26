# presentation/build_deck.py
# Generates PAINDICATOR.pptx (16:9) from the approved outline.
# Re-run any time to regenerate:  venv311\Scripts\python presentation\build_deck.py
# Requires: python-pptx, Pillow (both already in venv311).
# All comments in English only.

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
DOCS = os.path.join(ROOT, "docs")
SHOTS = os.path.join(DOCS, "screenshots")
OUT = os.path.join(HERE, "PAINDICATOR.pptx")

LOGO = os.path.join(DOCS, "logo.png")
IMG_WELCOME = os.path.join(SHOTS, "welcome.png")
IMG_FRONT = os.path.join(SHOTS, "painting-front.png")
IMG_BACK = os.path.join(SHOTS, "painting-back.png")
IMG_ANALYSIS = os.path.join(SHOTS, "clinician-analysis.png")

# ---------------------------------------------------------------------------
# Brand palette (project turquoise/teal from codes/theme + info.py accent)
# ---------------------------------------------------------------------------
ACCENT = RGBColor(0x00, 0xCE, 0xD1)      # turquoise
DARK = RGBColor(0x06, 0x3B, 0x3D)        # dark teal (headings / title bg)
DARK2 = RGBColor(0x0A, 0x53, 0x56)       # secondary dark teal
TEXT = RGBColor(0x26, 0x2B, 0x2B)        # near-black body text
MUTED = RGBColor(0x5A, 0x66, 0x66)       # muted captions
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE8, 0xF7, 0xF7)       # light teal tint (table rows / chips)
RED = RGBColor(0xC0, 0x39, 0x2B)         # red-flag accent

FONT = "Calibri"  # safe, ubiquitous; user can restyle later

# ---------------------------------------------------------------------------
# Presentation (16:9)
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(s, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill(shp, color)
    shp.shadow.inherit = False
    return shp


def textbox(s, left, top, width, height):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def set_run(run, text, size, color, bold=False, italic=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def para(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
         space_after=6, level=0, bullet=False, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    prefix = "•  " if bullet else ""
    set_run(p.add_run(), prefix + text, size, color, bold=bold, italic=italic)
    return p


def header(s, kicker, title):
    """Standard content-slide header: accent kicker + dark title + accent rule."""
    rect(s, 0, 0, SW, Inches(0.18), ACCENT)  # top accent strip
    tb, tf = textbox(s, Inches(0.7), Inches(0.45), Inches(12.0), Inches(1.4))
    para(tf, kicker.upper(), 14, ACCENT, bold=True, space_after=2, first=True)
    para(tf, title, 32, DARK, bold=True, space_after=0)
    rect(s, Inches(0.72), Inches(1.72), Inches(2.2), Inches(0.06), ACCENT)


def footer(s, idx):
    tb, tf = textbox(s, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4))
    p = para(tf, "PAINDICATOR  ·  Tel Aviv University × Loewenstein (Clalit)",
             10, MUTED, first=True)
    # page number on the right
    tb2, tf2 = textbox(s, Inches(12.3), Inches(7.0), Inches(0.9), Inches(0.4))
    para(tf2, str(idx), 10, MUTED, align=PP_ALIGN.RIGHT, first=True)


def add_image_contain(s, path, left, top, width, height, border=True):
    """Place an image fully inside the box, preserving aspect ratio, centered."""
    if not os.path.exists(path):
        # Visible placeholder so the deck still builds if an asset is missing.
        ph = rect(s, left, top, width, height, LIGHT)
        tb, tf = textbox(s, left, top + height / 2 - Inches(0.2), width, Inches(0.4))
        para(tf, "[image missing: %s]" % os.path.basename(path), 11, MUTED,
             align=PP_ALIGN.CENTER, first=True)
        return ph
    iw, ih = Image.open(path).size
    box_w, box_h = int(width), int(height)
    scale = min(box_w / iw, box_h / ih)
    new_w, new_h = int(iw * scale), int(ih * scale)
    nl = int(left) + (box_w - new_w) // 2
    nt = int(top) + (box_h - new_h) // 2
    if border:
        rect(s, Emu(nl - 9525), Emu(nt - 9525), Emu(new_w + 19050), Emu(new_h + 19050), ACCENT)
    return s.shapes.add_picture(path, Emu(nl), Emu(nt), width=Emu(new_w), height=Emu(new_h))


def chip(s, left, top, text, color=DARK2):
    w = Inches(0.18 * len(text) + 0.5)
    shp = rect(s, left, top, w, Inches(0.42), color)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), text, 12, WHITE, bold=True)
    return left + w + Inches(0.18)


# ===========================================================================
# SLIDE 1 — Title & team
# ===========================================================================
s = slide()
rect(s, 0, 0, SW, SH, DARK)                       # full dark background
rect(s, 0, 0, SW, Inches(0.25), ACCENT)
rect(s, 0, SH - Inches(0.25), SW, Inches(0.25), ACCENT)
add_image_contain(s, LOGO, Inches(4.17), Inches(0.9), Inches(5.0), Inches(2.4), border=False)

tb, tf = textbox(s, Inches(1.0), Inches(3.5), Inches(11.33), Inches(2.6))
para(tf, "PAINDICATOR", 46, WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=2, first=True)
para(tf, "Map Your Pain", 22, ACCENT, bold=True, italic=True, align=PP_ALIGN.CENTER, space_after=14)
para(tf, "A clinical decision-support system for digital pain mapping & dermatome analysis",
     18, LIGHT, align=PP_ALIGN.CENTER, space_after=18)
para(tf, "Amit Bezalel  ·  Shachar Guttman   |   Biomedical Engineering, Tel Aviv University",
     16, WHITE, align=PP_ALIGN.CENTER, space_after=4)
para(tf, "Supervisor: Prof. Mickey Scheinowitz   ·   In collaboration with Loewenstein "
         "Rehabilitation Medical Center (Clalit Health Services)",
     13, MUTED if False else LIGHT, align=PP_ALIGN.CENTER, space_after=0)

# ===========================================================================
# SLIDE 2 — The problem
# ===========================================================================
s = slide()
header(s, "The problem", "Pain is everywhere — and barely documented")
tb, tf = textbox(s, Inches(0.7), Inches(2.2), Inches(7.2), Inches(4.4))
para(tf, "Pain is a top reason people see a doctor — yet one of the worst-documented "
         "clinical signals.", 18, TEXT, space_after=14, first=True)
for t in [
    "Recorded as free text or a felt-tip drawing on a paper body chart",
    "Subjective & unstructured — no numbers to compute",
    "Not comparable over time — pain care is about change, and we throw it away",
    "Loses anatomical precision — a flat sketch doesn't point to a spinal level",
]:
    para(tf, t, 16, TEXT, bullet=True, space_after=10)
para(tf, "The clinician's real job — localizing pain to a spinal nerve level — is done "
         "from memory and a 2D picture.", 16, DARK2, bold=True, space_after=0)
# right panel
rect(s, Inches(8.3), Inches(2.2), Inches(4.3), Inches(4.0), LIGHT)
tb, tf = textbox(s, Inches(8.5), Inches(2.6), Inches(3.9), Inches(3.2))
para(tf, "TODAY", 14, ACCENT, bold=True, align=PP_ALIGN.CENTER, space_after=8, first=True)
para(tf, "✎  paper drawing", 20, DARK, bold=True, align=PP_ALIGN.CENTER, space_after=10)
para(tf, "“lower back, radiating\ndown the leg”", 18, MUTED, italic=True,
     align=PP_ALIGN.CENTER, space_after=10)
para(tf, "→ subjective\n→ not measurable\n→ not trackable", 16, RED, bold=True,
     align=PP_ALIGN.CENTER, space_after=0)
footer(s, 2)

# ===========================================================================
# SLIDE 3 — The insight
# ===========================================================================
s = slide()
header(s, "The insight", "Paint pain on a 3D body — every point is a nerve")
add_image_contain(s, IMG_FRONT, Inches(7.2), Inches(2.05), Inches(5.6), Inches(4.7))
tb, tf = textbox(s, Inches(0.7), Inches(2.3), Inches(6.3), Inches(4.3))
para(tf, "What if the patient painted their pain onto an interactive 3D model?", 18, TEXT,
     bold=True, space_after=14, first=True)
for t in [
    "Three intensity levels — mild / moderate / severe",
    "Every vertex is mapped to a dermatome (the skin of one spinal nerve root)",
    "So painting isn't a picture — it becomes anatomical data, automatically",
]:
    para(tf, t, 16, TEXT, bullet=True, space_after=12)
para(tf, "Subjective report  →  quantitative, anatomical signal", 18, ACCENT, bold=True,
     space_after=0)
footer(s, 3)

# ===========================================================================
# SLIDE 4 — Demo / walkthrough
# ===========================================================================
s = slide()
header(s, "See it work", "From patient drawing to clinician analysis")
# 2x2 screenshot grid
gw, gh = Inches(5.7), Inches(2.15)
add_image_contain(s, IMG_WELCOME, Inches(0.7), Inches(2.15), gw, gh)
add_image_contain(s, IMG_FRONT, Inches(6.9), Inches(2.15), gw, gh)
add_image_contain(s, IMG_BACK, Inches(0.7), Inches(4.45), gw, gh)
add_image_contain(s, IMG_ANALYSIS, Inches(6.9), Inches(4.45), gw, gh)
# caption chips (dark, readable) at the top-left of each cell
caps = [
    (Inches(0.7), Inches(2.15), "1 · Guided entry"),
    (Inches(6.9), Inches(2.15), "2 · Paint the front"),
    (Inches(0.7), Inches(4.45), "3 · Rotate & paint the back"),
    (Inches(6.9), Inches(4.45), "4 · Clinician analysis"),
]
for cx, cy, cap in caps:
    chip(s, cx, cy, cap, color=DARK)
tb, tf = textbox(s, Inches(0.7), Inches(6.75), Inches(12.0), Inches(0.4))
para(tf, "Live demo optional — these captured screens are the fallback flow.", 12, MUTED,
     italic=True, first=True)
footer(s, 4)

# ===========================================================================
# SLIDE 5 — From paint to data (pipeline)
# ===========================================================================
s = slide()
header(s, "How it works", "From paint to ranked, per-nerve metrics")
steps = [
    ("Patient paints", "per-vertex levels 0–3 on the 3D model"),
    ("Dermatome lookup", "1 byte / vertex map → C1–S5 (31 dermatomes)"),
    ("Coverage engine", "pure math · no graphics · fully testable"),
    ("Ranked output", "per-dermatome metrics + clinician panel + saved JSON"),
]
y = Inches(2.4)
for i, (t, d) in enumerate(steps):
    box = rect(s, Inches(0.9), y, Inches(11.5), Inches(0.92), LIGHT if i % 2 == 0 else WHITE)
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1.25)
    num = rect(s, Inches(1.1), y + Inches(0.16), Inches(0.6), Inches(0.6), ACCENT)
    ntf = num.text_frame
    ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    np_ = ntf.paragraphs[0]
    np_.alignment = PP_ALIGN.CENTER
    set_run(np_.add_run(), str(i + 1), 20, WHITE, bold=True)
    tb, tf = textbox(s, Inches(2.0), y + Inches(0.08), Inches(10.2), Inches(0.78))
    para(tf, t, 18, DARK, bold=True, space_after=0, first=True)
    para(tf, d, 14, MUTED, space_after=0)
    y = y + Inches(1.06)
footer(s, 5)

# ===========================================================================
# SLIDE 6 — Quantitative analytics (metrics table)
# ===========================================================================
s = slide()
header(s, "Clinical value", "Quantitative dermatome analytics")
add_image_contain(s, IMG_ANALYSIS, Inches(8.5), Inches(2.05), Inches(4.2), Inches(3.6))
metrics = [
    ("Metric", "What it tells the clinician"),
    ("Weighted pain burden", "this nerve level's share of total pain → ranks the dominant level"),
    ("Pain area share", "large mild area vs small intense area"),
    ("Local involvement", "how much of the nerve zone is affected"),
    ("Mean local intensity", "average severity where painted (1–3)"),
    ("Segmental spread", "narrow = single level · wide = central/multifocal"),
    ("Overlap flag", "adjacent levels near-equal → radiculopathy ambiguity"),
]
rows, cols = len(metrics), 2
tbl_shape = s.shapes.add_table(rows, cols, Inches(0.7), Inches(2.05),
                               Inches(7.5), Inches(4.2))
table = tbl_shape.table
table.columns[0].width = Inches(2.7)
table.columns[1].width = Inches(4.8)
for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        txt = metrics[r][c]
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = DARK
            set_run(p.add_run(), txt, 13, WHITE, bold=True)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if r % 2 else WHITE
            set_run(p.add_run(), txt, 12, TEXT if c else DARK2,
                    bold=(c == 0))
tb, tf = textbox(s, Inches(8.5), Inches(5.75), Inches(4.3), Inches(1.2))
para(tf, "“L5–S1 distribution, 78% burden in L5, sharp 3/3 → consistent with L5–S1 disc "
         "herniation.”", 13, DARK, bold=True, italic=True, first=True)
footer(s, 6)

# ===========================================================================
# SLIDE 7 — Decision support
# ===========================================================================
s = slide()
header(s, "Clinical + technical", "Explainable decision support")
tb, tf = textbox(s, Inches(0.7), Inches(2.2), Inches(6.2), Inches(4.4))
para(tf, "A rule-based analyzer correlates the dermatome pattern with the questionnaire:",
     17, TEXT, bold=True, space_after=12, first=True)
for t in [
    "Localizing patterns — e.g. L5–S1 sciatica, C6–C7 disc, femoral nerve (L2–L4)",
    "Symptom phenotyping — burning+tingling → neuropathic; Valsalva-positive → disc",
    "Widespread (≥7 levels) → central sensitization / fibromyalgia",
]:
    para(tf, t, 15, TEXT, bullet=True, space_after=10)
para(tf, "Every line is a readable rule a clinician can audit & override — explainable, "
         "not a black box.", 15, DARK2, bold=True, space_after=0)
# red-flags panel
rect(s, Inches(7.2), Inches(2.2), Inches(5.4), Inches(3.0), RGBColor(0xFB, 0xEC, 0xEA))
rect(s, Inches(7.2), Inches(2.2), Inches(0.12), Inches(3.0), RED)
tb, tf = textbox(s, Inches(7.55), Inches(2.4), Inches(5.0), Inches(2.7))
para(tf, "RED FLAGS surfaced", 15, RED, bold=True, space_after=8, first=True)
for t in ["Cauda equina (S3–S4–S5 / bladder–bowel) → urgent referral",
          "Severe unilateral thoracic band → herpes zoster (pre-rash)",
          "Fever + spine pain → discitis / epidural abscess",
          "Multi-level cervical/thoracic → myelopathy"]:
    para(tf, t, 13, TEXT, bullet=True, space_after=6)
rect(s, Inches(7.2), Inches(5.45), Inches(5.4), Inches(1.1), DARK)
tb, tf = textbox(s, Inches(7.45), Inches(5.6), Inches(5.0), Inches(0.85))
para(tf, "Decision support — descriptive only. Not a diagnosis; not a substitute for "
         "clinical judgment.", 13, WHITE, bold=True, first=True)
footer(s, 7)

# ===========================================================================
# SLIDE 8 — Under the hood
# ===========================================================================
s = slide()
header(s, "Engineering", "A real, robust application")
tb, tf = textbox(s, Inches(0.7), Inches(2.15), Inches(7.0), Inches(4.4))
para(tf, "~5,600 lines of Python · clean ownership split: rendering / interaction / "
         "analytics core", 16, TEXT, bold=True, space_after=12, first=True)
for t in [
    "PyQt6 UI + VTK 3D pipeline (lazy init, crash-safe lifecycle)",
    "Per-vertex radius brush with normal gating — can't paint through the body",
    "Gap-free fast strokes + true per-stroke undo + erase",
    "Tablet-first: stylus pressure, two-finger pinch / pan / rotate, gesture guide",
    "Bilingual English / Hebrew with full RTL",
    "Ships as a standalone, offline Windows EXE (PyInstaller)",
]:
    para(tf, t, 15, TEXT, bullet=True, space_after=9)
# chips column
tb, tf = textbox(s, Inches(8.1), Inches(2.15), Inches(4.5), Inches(0.5))
para(tf, "STACK", 14, ACCENT, bold=True, first=True)
y = Inches(2.7)
x = Inches(8.1)
for label in ["Python 3.11", "PyQt6 6.10", "VTK 9.5.2", "NumPy 2.3.5",
              "Pillow 12.0", "PyInstaller 6.17"]:
    chip(s, x, y, label)
    y = y + Inches(0.6)
# scope box
rect(s, Inches(8.1), Inches(6.0), Inches(4.5), Inches(0.7), LIGHT)
tb, tf = textbox(s, Inches(8.25), Inches(6.08), Inches(4.2), Inches(0.55))
para(tf, "2 models · 31 dermatomes · 12 screens", 14, DARK, bold=True,
     align=PP_ALIGN.CENTER, first=True)
footer(s, 8)

# ===========================================================================
# SLIDE 9 — Impact & partner
# ===========================================================================
s = slide()
header(s, "Real-world impact", "Built with a clinical partner")
tb, tf = textbox(s, Inches(0.7), Inches(2.2), Inches(7.2), Inches(4.2))
para(tf, "Co-developed with Loewenstein Rehabilitation Medical Center (Clalit Health "
         "Services) around a real documentation problem.", 18, TEXT, bold=True,
     space_after=16, first=True)
para(tf, "Where it helps:", 16, DARK2, bold=True, space_after=8)
for t in ["Pain clinics", "Neurology", "Orthopedics & spine",
          "Physical medicine & rehabilitation"]:
    para(tf, t, 16, TEXT, bullet=True, space_after=8)
para(tf, "Sessions are saved and comparable over time — we finally measure how pain "
         "changes.", 16, ACCENT, bold=True, space_after=0)
# partner panel
rect(s, Inches(8.4), Inches(2.2), Inches(4.2), Inches(4.0), DARK)
tb, tf = textbox(s, Inches(8.7), Inches(2.6), Inches(3.6), Inches(3.4))
para(tf, "PARTNERS", 14, ACCENT, bold=True, align=PP_ALIGN.CENTER, space_after=14, first=True)
para(tf, "Tel Aviv University", 20, WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=4)
para(tf, "Biomedical Engineering", 13, LIGHT, align=PP_ALIGN.CENTER, space_after=18)
para(tf, "Loewenstein Rehabilitation\nMedical Center", 18, WHITE, bold=True,
     align=PP_ALIGN.CENTER, space_after=4)
para(tf, "Clalit Health Services", 13, LIGHT, align=PP_ALIGN.CENTER, space_after=0)
footer(s, 9)

# ===========================================================================
# SLIDE 10 — Roadmap
# ===========================================================================
s = slide()
header(s, "Vision", "Where it goes next")
items = [
    ("Visual heatmap", "color the body by computed pain burden — at a glance"),
    ("Longitudinal metrics", "numeric change-over-time from the multi-session overlay"),
    ("ML pattern recognition", "rules become the labeled foundation + safety net"),
    ("Regulatory path", "position as decision support; explore CE / device class"),
]
x = Inches(0.7)
cw = Inches(2.95)
for i, (t, d) in enumerate(items):
    left = Inches(0.7 + i * 3.07)
    rect(s, left, Inches(2.5), cw, Inches(3.2), LIGHT)
    rect(s, left, Inches(2.5), cw, Inches(0.12), ACCENT)
    num = rect(s, left + Inches(0.3), Inches(2.85), Inches(0.7), Inches(0.7), ACCENT)
    ntf = num.text_frame; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    npp = ntf.paragraphs[0]; npp.alignment = PP_ALIGN.CENTER
    set_run(npp.add_run(), str(i + 1), 22, WHITE, bold=True)
    tb, tf = textbox(s, left + Inches(0.2), Inches(3.8), cw - Inches(0.4), Inches(1.8))
    para(tf, t, 16, DARK, bold=True, space_after=8, first=True)
    para(tf, d, 13, MUTED, space_after=0)
footer(s, 10)

# ===========================================================================
# SLIDE 11 — Why we win
# ===========================================================================
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, 0, SW, Inches(0.2), ACCENT)
tb, tf = textbox(s, Inches(0.8), Inches(0.55), Inches(11.7), Inches(1.0))
para(tf, "WHY WE WIN", 16, ACCENT, bold=True, space_after=2, first=True)
para(tf, "Five reasons to back PAINDICATOR", 32, WHITE, bold=True, space_after=0)
points = [
    ("Genuine novelty", "quantitative pain → nerve mapping, ranked burden + overlap flag"),
    ("Real clinical partner", "co-developed with Loewenstein / Clalit"),
    ("Finished & deployable", "offline tablet app — works today, end to end"),
    ("Explainable & private", "auditable rules; local data; no cloud"),
    ("Solid engineering", "~5,600 LOC · 2 models · 31 dermatomes"),
]
y = Inches(1.95)
for i, (t, d) in enumerate(points):
    num = rect(s, Inches(0.9), y, Inches(0.62), Inches(0.62), ACCENT)
    ntf = num.text_frame; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    npp = ntf.paragraphs[0]; npp.alignment = PP_ALIGN.CENTER
    set_run(npp.add_run(), str(i + 1), 20, DARK, bold=True)
    tb, tf = textbox(s, Inches(1.75), y - Inches(0.05), Inches(10.8), Inches(0.95))
    para(tf, t, 20, WHITE, bold=True, space_after=0, first=True)
    para(tf, d, 14, LIGHT, space_after=0)
    y = y + Inches(0.98)
footer(s, 11)

# ===========================================================================
# SLIDE 12 — Thank you / disclaimer
# ===========================================================================
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, 0, SW, Inches(0.25), ACCENT)
rect(s, 0, SH - Inches(0.25), SW, Inches(0.25), ACCENT)
add_image_contain(s, LOGO, Inches(4.67), Inches(0.95), Inches(4.0), Inches(1.9), border=False)
tb, tf = textbox(s, Inches(1.0), Inches(3.1), Inches(11.33), Inches(2.6))
para(tf, "Thank you", 40, WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=8, first=True)
para(tf, "Amit Bezalel  ·  Shachar Guttman   —   Biomedical Engineering, Tel Aviv University",
     16, LIGHT, align=PP_ALIGN.CENTER, space_after=4)
para(tf, "We'd love your questions.", 18, ACCENT, bold=True, align=PP_ALIGN.CENTER, space_after=0)
# disclaimer band
rect(s, Inches(1.5), Inches(5.7), Inches(10.33), Inches(1.1), DARK2)
tb, tf = textbox(s, Inches(1.8), Inches(5.85), Inches(9.7), Inches(0.85))
para(tf, "PAINDICATOR is a research / decision-support tool. Its analysis is a descriptive "
         "summary of recorded data — not a diagnosis, and not a substitute for clinical "
         "judgment, physical examination, or imaging.", 12, WHITE, italic=True,
     align=PP_ALIGN.CENTER, first=True)

# ---------------------------------------------------------------------------
prs.save(OUT)
print("Wrote %s (%d slides)" % (OUT, len(prs.slides._sldIdLst)))
